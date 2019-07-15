"""
    Author:         <REPLACE>
    Project:        100DaysPython
    File:           module3_day34_powerpoint.py
    Creation Date:  <REPLACE>
    Description:    <REPLACE>
"""

from pptx import Presentation
from datetime import datetime

pres = "TheDoctor.pptx"
doctors = Presentation(pres)
print(f"doctors: {doctors}")

for i, slide in enumerate(doctors.slides):
    print(f"{i}: {slide}")

title_slide, content_slide = doctors.slides[0], doctors.slides[1]
print(f"Slide ID: {title_slide.slide_id}\nSlide Element: {title_slide.element}\nSlide Layout: "
      f"{title_slide.slide_layout.name}")
print(f"Slide ID: {content_slide.slide_id}\nSlide Element: {content_slide.element}\nSlide Layout: "
      f"{content_slide.slide_layout.name}")

for i, slide in enumerate(doctors.slides, start=1):
    print(f"Slide {i}, Type: {slide.slide_layout.name}")
    for shape in slide.shapes:
        print(f"Shape Type: {shape.placeholder_format.idx}\tPlaceholder Name: {shape.name}")

doctor = dict()
for i, slide in enumerate(doctors.slides, start=1):
    if slide.slide_layout.name == "Two Content":
        title = slide.shapes.title.text
        deets_dict = dict()
        for shape in slide.shapes:
            if shape.placeholder_format.idx == 1:
                details = shape.text.split("\n")
                for item in details:
                    deets_dict[item.split(": ")[0]] = item.split(": ")[1]
                doctor[title] = deets_dict
    else:
        continue

print(doctor)

# A blank presentation object is created.
new_pres = Presentation()

# The title slide is added which includes a creation date in the subtitle.
slide = new_pres.slides.add_slide(new_pres.slide_layouts[0])
slide.shapes.title.text = "This was created with python".title()
slide.placeholders[1].text = f"Created on {datetime.date(datetime.now())}".title()

# The contents of the dictionary are iterated over to create a new slide for each new doctor.
for i, title in enumerate(doctor):
    new_slide = new_pres.slides.add_slide(new_pres.slide_layouts[3])
    shapes = new_slide.shapes
    title_shape = shapes.title
    title_shape.text = title
    # The details of the Doctor are added to the left placeholder. The actor is added first and then the remaining
    # details are iterated over to add as a new paragraph.
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = f"{list(doctor[title].keys())[0]}: {doctor[title][list(doctor[title].keys())[0]]}"
    for j in range(1, len(doctor[title])):
        tf.add_paragraph().text = f"{list(doctor[title].keys())[j]}: {doctor[title][list(doctor[title].keys())[j]]}"

# After the content is added, the presentation is saved.
new_pres.save('auto_pres.pptx')